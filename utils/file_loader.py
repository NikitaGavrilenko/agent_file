import os
import asyncio
from pathlib import Path
from typing import List, Optional
from loguru import logger

from models.document import Document


class FileLoader:
    """Класс для загрузки и обработки различных типов файлов"""
    
    def __init__(self, supported_formats: List[str] = None):
        if supported_formats is None:
            supported_formats = [".pdf", ".docx", ".pptx", ".txt"]
        self.supported_formats = supported_formats
    
    async def load_documents_from_folder(self, folder_path: str) -> List[Document]:
        """Загружает все документы из папки"""
        documents = []
        folder = Path(folder_path)
        
        if not folder.exists():
            raise FileNotFoundError(f"Папка не найдена: {folder_path}")
        
        if not folder.is_dir():
            raise NotADirectoryError(f"Путь не является папкой: {folder_path}")
        
        # Получаем список всех файлов
        files = [f for f in folder.iterdir() if f.is_file()]
        
        # Фильтруем по поддерживаемым форматам
        supported_files = [
            f for f in files 
            if f.suffix.lower() in self.supported_formats
        ]
        
        if not supported_files:
            logger.warning(f"В папке {folder_path} не найдено поддерживаемых файлов")
            return []
        
        logger.info(f"Найдено {len(supported_files)} поддерживаемых файлов")
        
        # Загружаем каждый файл
        for file_path in supported_files:
            try:
                document = await self.load_single_document(file_path)
                if document:
                    documents.append(document)
            except Exception as e:
                logger.error(f"Ошибка при загрузке файла {file_path}: {e}")
                continue
        
        logger.info(f"Успешно загружено {len(documents)} документов")
        return documents
    
    async def load_single_document(self, file_path: Path) -> Optional[Document]:
        """Загружает один документ"""
        try:
            file_type = file_path.suffix.lower()
            
            if file_type == ".pdf":
                content = await self._load_pdf(file_path)
            elif file_type == ".docx":
                content = await self._load_docx(file_path)
            elif file_type == ".pptx":
                content = await self._load_pptx(file_path)
            elif file_type == ".txt":
                content = await self._load_txt(file_path)
            else:
                logger.warning(f"Неподдерживаемый тип файла: {file_type}")
                return None
            
            if not content:
                logger.warning(f"Не удалось извлечь содержимое из файла: {file_path}")
                return None
            
            # Очищаем и обрезаем текст
            content = self._clean_content(content)
            
            return Document.from_file(str(file_path), content)
            
        except Exception as e:
            logger.error(f"Ошибка при загрузке файла {file_path}: {e}")
            return None
    
    async def _load_pdf(self, file_path: Path) -> str:
        """Загружает содержимое PDF файла"""
        try:
            import PyPDF2
            
            content = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    content += page.extract_text() + "\n"
            
            return content
            
        except ImportError:
            logger.error("PyPDF2 не установлен. Установите: pip install PyPDF2")
            return ""
        except Exception as e:
            logger.error(f"Ошибка при чтении PDF {file_path}: {e}")
            return ""
    
    async def _load_docx(self, file_path: Path) -> str:
        """Загружает содержимое Word файла"""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(file_path)
            content = ""
            
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            
            # Добавляем содержимое таблиц
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        content += cell.text + "\t"
                    content += "\n"
            
            return content
            
        except ImportError:
            logger.error("python-docx не установлен. Установите: pip install python-docx")
            return ""
        except Exception as e:
            logger.error(f"Ошибка при чтении DOCX {file_path}: {e}")
            return ""
    
    async def _load_pptx(self, file_path: Path) -> str:
        """Загружает содержимое PowerPoint файла"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            
            return content
            
        except ImportError:
            logger.error("python-pptx не установлен. Установите: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"Ошибка при чтении PPTX {file_path}: {e}")
            return ""
    
    async def _load_txt(self, file_path: Path) -> str:
        """Загружает содержимое текстового файла"""
        try:
            # Пробуем разные кодировки
            encodings = ['utf-8', 'cp1251', 'windows-1251', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            
            # Если не удалось с кодировками, пробуем бинарный режим
            with open(file_path, 'rb') as file:
                return file.read().decode('utf-8', errors='ignore')
                
        except Exception as e:
            logger.error(f"Ошибка при чтении TXT {file_path}: {e}")
            return ""
    
    def _clean_content(self, content: str) -> str:
        """Очищает содержимое от лишних символов"""
        if not content:
            return ""
        
        # Убираем множественные пробелы
        import re
        content = re.sub(r'\s+', ' ', content)
        
        # Убираем множественные переносы строк
        content = re.sub(r'\n+', '\n', content)
        
        # Убираем лишние пробелы в начале и конце
        content = content.strip()
        
        return content
    
    def get_file_info(self, file_path: Path) -> dict:
        """Получает информацию о файле"""
        try:
            stat = file_path.stat()
            return {
                "name": file_path.name,
                "size": stat.st_size,
                "created": stat.st_ctime,
                "modified": stat.st_mtime,
                "type": file_path.suffix.lower()
            }
        except Exception as e:
            logger.error(f"Ошибка при получении информации о файле {file_path}: {e}")
            return {}
